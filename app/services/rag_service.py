import os
import pdfplumber
import pandas as pd

from docx import Document
from pptx import Presentation
from openai import OpenAI
from sentence_transformers import SentenceTransformer
from sqlalchemy.orm import Session

from app.core.config import settings
from app.models.document import DocumentChunk
from app.schemas.document import AskResponse, ChunkResult, SearchResponse

# Load embedding model once
_model = SentenceTransformer("all-MiniLM-L6-v2")


# --------------------------------------------------
# TEXT EXTRACTION FUNCTIONS
# --------------------------------------------------

def _extract_pdf_text(path: str) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
    return "\n".join(pages)


def _extract_docx_text(path: str) -> str:
    doc = Document(path)
    return "\n".join(
        paragraph.text
        for paragraph in doc.paragraphs
        if paragraph.text.strip()
    )


def _extract_pptx_text(path: str) -> str:
    prs = Presentation(path)

    slides = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    slides.append(shape.text)

    return "\n".join(slides)


def _extract_txt_text(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def _extract_csv_text(path: str) -> str:
    df = pd.read_csv(path)
    return df.astype(str).to_string(index=False)


# --------------------------------------------------
# COMMON EXTRACTOR
# --------------------------------------------------

def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return _extract_pdf_text(file_path)

    elif ext == ".docx":
        return _extract_docx_text(file_path)

    elif ext == ".pptx":
        return _extract_pptx_text(file_path)

    elif ext == ".txt":
        return _extract_txt_text(file_path)

    elif ext == ".csv":
        return _extract_csv_text(file_path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


# --------------------------------------------------
# CHUNKING
# --------------------------------------------------

def _chunk_text(
    content: str,
    chunk_size: int = 500,
    overlap: int = 100,
) -> list[str]:

    chunks = []
    start = 0

    while start < len(content):
        end = start + chunk_size
        chunks.append(content[start:end])
        start += chunk_size - overlap

    return chunks


# --------------------------------------------------
# UPLOAD DOCUMENT
# --------------------------------------------------

def upload_document(
    db: Session,
    filename: str,
    file_path: str,
) -> int:

    raw_text = extract_text(file_path)

    if not raw_text.strip():
        raise ValueError("No text could be extracted from the document.")

    chunks = _chunk_text(raw_text)

    embeddings = _model.encode(
        chunks,
        show_progress_bar=False,
    )

    db.query(DocumentChunk).filter(
        DocumentChunk.filename == filename
    ).delete()

    for chunk_text, embedding in zip(chunks, embeddings):
        db.add(
            DocumentChunk(
                filename=filename,
                text=chunk_text,
                embedding=embedding.tolist(),
            )
        )

    db.commit()

    return len(chunks)


# --------------------------------------------------
# SEARCH
# --------------------------------------------------

def search(
    db: Session,
    query: str,
    top_k: int = 3,
) -> SearchResponse:

    query_embedding = _model.encode([query])[0].tolist()

    rows = (
        db.query(
            DocumentChunk.text,
            DocumentChunk.embedding.cosine_distance(query_embedding).label("distance"),
        )
        .order_by("distance")
        .limit(top_k)
        .all()
    )

    results = [
        ChunkResult(
            text=row.text,
            score=round(1 - row.distance, 4),
        )
        for row in rows
    ]

    return SearchResponse(results=results)


# --------------------------------------------------
# ASK LLM
# --------------------------------------------------

def ask(
    db: Session,
    query: str,
    top_k: int = 3,
) -> AskResponse:

    search_result = search(db, query, top_k)

    context_text = "\n\n".join(
        chunk.text
        for chunk in search_result.results
    )

    client = OpenAI(
        api_key=settings.GROQ_API_KEY,
        base_url="https://api.groq.com/openai/v1",
    )

    prompt = f"""
You are a helpful assistant.

Answer ONLY using the provided context.

Question:
{query}

Context:
{context_text}

If the answer is not found in the context, reply:
"I couldn't find that information in the uploaded documents."
"""

    response = client.chat.completions.create(
        model=settings.GROQ_MODEL,
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
    )

    return AskResponse(
        answer=response.choices[0].message.content,
        context=search_result.results,
    )